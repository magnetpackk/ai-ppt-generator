# -*- coding: utf-8 -*-
"""
模块一：Template Parser - PPT模板结构解析
解析PPT模板，识别每页类型和占位区域
"""

import os
import io
import base64
from typing import List, Dict, Optional, Tuple
from dataclasses import dataclass, field
from io import BytesIO

from pptx import Presentation
from pptx.util import Emu, Pt
import openai


@dataclass
class Placeholder:
    """占位区域信息"""
    shape_id: int
    type: str  # text | image | chart | table | diagram
    x: float
    y: float
    width: float
    height: float
    expected_text_length: Optional[int] = None

    @property
    def bounding_box(self) -> Dict[str, float]:
        return {
            "x": self.x,
            "y": self.y,
            "w": self.width,
            "h": self.height,
        }


@dataclass
class SlideInfo:
    """单页幻灯片信息"""
    index: int
    page_type: str  # cover | toc | title | content | section | blank
    placeholders: List[Placeholder]
    text_areas: List[str] = field(default_factory=list)


@dataclass
class ThemeInfo:
    """主题信息"""
    background_color: Optional[str] = None
    font_family: Optional[str] = None
    color_scheme: List[str] = field(default_factory=list)


@dataclass
class TemplateStructure:
    """模板整体结构"""
    slide_count: int
    slides: List[SlideInfo]
    theme: ThemeInfo


class TemplateParser:
    """PPT模板解析器"""

    def __init__(self, openai_api_key: str, openai_base_url: Optional[str] = None):
        self.openai_api_key = openai_api_key
        self.client = openai.OpenAI(
            api_key=openai_api_key,
            base_url=openai_base_url,
        )

    def parse(self, template_file: BytesIO) -> TemplateStructure:
        """解析PPT模板文件，返回结构化的模板结构"""
        prs = Presentation(template_file)

        slides_info: List[SlideInfo] = []
        theme = self._extract_theme(prs)

        for idx, slide in enumerate(prs.slides):
            slide_info = self._analyze_slide(idx, slide)
            slides_info.append(slide_info)

        return TemplateStructure(
            slide_count=len(prs.slides),
            slides=slides_info,
            theme=theme,
        )

    def _extract_theme(self, prs: Presentation) -> ThemeInfo:
        """提取主题信息"""
        theme = ThemeInfo()

        # 尝试提取主题配色
        try:
            if prs.slide_masters and len(prs.slide_masters) > 0:
                master = prs.slide_masters[0]
                # 简单提取，实际需要更复杂的处理
                theme.color_scheme = []
                for shape in master.shapes:
                    try:
                        if hasattr(shape, 'fill') and shape.fill:
                            color = shape.fill.fore_color
                            if color.type == "rgb":
                                rgb = color.rgb
                                if rgb:
                                    hex_color = f"#{rgb}"
                                    theme.color_scheme.append(hex_color)
                    except Exception:
                        pass
        except Exception:
            pass

        return theme

    def _analyze_slide(self, index: int, slide) -> SlideInfo:
        """分析单页幻灯片，识别页面类型和占位区域"""
        placeholders: List[Placeholder] = []
        text_areas: List[str] = []

        for shape in slide.shapes:
            ph = self._process_shape(shape)
            if ph:
                placeholders.append(ph)

            # 收集文本内容用于AI判断页面类型
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        text_areas.append(text)

        # 使用AI判断页面类型和占位区域用途
        page_type = self._ai_classify_page(index, text_areas, placeholders)

        return SlideInfo(
            index=index,
            page_type=page_type,
            placeholders=placeholders,
            text_areas=text_areas,
        )

    def _process_shape(self, shape) -> Optional[Placeholder]:
        """处理单个形状，判断是否是占位区域"""
        # 获取形状位置（单位：英寸，方便AI理解）
        x = shape.left / 914400  # EMU to inches
        y = shape.top / 914400
        w = shape.width / 914400
        h = shape.height / 914400

        shape_id = shape.shape_id

        # 判断形状类型
        if shape.has_text_frame:
            # 统计文本长度，预估适合放多少文字
            total_text = ""
            for p in shape.text_frame.paragraphs:
                total_text += p.text.strip() + " "
            text_len = len(total_text)

            return Placeholder(
                shape_id=shape_id,
                type="text",
                x=round(x, 2),
                y=round(y, 2),
                width=round(w, 2),
                height=round(h, 2),
                expected_text_length=text_len if text_len > 0 else None,
            )

        elif shape.is_placeholder:
            ph_type = shape.placeholder_format.type
            if ph_type in [1, 14]:  # Title or Body
                return Placeholder(
                    shape_id=shape_id,
                    type="text",
                    x=round(x, 2),
                    y=round(y, 2),
                    width=round(w, 2),
                    height=round(h, 2),
                )
            elif ph_type == 18:  # Picture
                return Placeholder(
                    shape_id=shape_id,
                    type="image",
                    x=round(x, 2),
                    y=round(y, 2),
                    width=round(w, 2),
                    height=round(h, 2),
                )

        elif shape.has_chart:
            return Placeholder(
                shape_id=shape_id,
                type="chart",
                x=round(x, 2),
                y=round(y, 2),
                width=round(w, 2),
                height=round(h, 2),
            )

        elif shape.has_table:
            return Placeholder(
                shape_id=shape_id,
                type="table",
                x=round(x, 2),
                y=round(y, 2),
                width=round(w, 2),
                height=round(h, 2),
            )

        # 判断是否是图片占位框（空的矩形）
        try:
            if not shape.has_text_frame and (shape.fill is None or (hasattr(shape, 'fill') and shape.fill.type is None)):
                # 空矩形很可能是图片占位
                if w > h * 1.2 or h > w * 1.2:  # 明显的图片比例
                    return Placeholder(
                        shape_id=shape_id,
                        type="image",
                        x=round(x, 2),
                        y=round(y, 2),
                        width=round(w, 2),
                        height=round(h, 2),
                    )
        except AttributeError:
            # GroupShape 等特殊形状可能没有 fill 属性，直接跳过
            pass

        return None

    def _ai_classify_page(self, index: int, text_areas: List[str], placeholders: List[Placeholder]) -> str:
        """使用AI分类页面类型"""

        # 如果已有明确文字，直接判断
        combined_text = " ".join(text_areas).lower()
        if "目录" in combined_text or "contents" in combined_text or "toc" in combined_text:
            return "toc"
        if "封面" in combined_text or "title" in combined_text:
            return "cover"
        if len(placeholders) == 0:
            return "blank"
        if len(text_areas) <= 2 and any(p.type == "text" for p in placeholders) and len(placeholders) <= 2:
            return "title"
        if "章节" in combined_text or "section" in combined_text:
            return "section"

        # 无法通过规则判断，调用LLM
        prompt = self._build_classification_prompt(index, text_areas, placeholders)

        try:
            response = self.client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {
                        "role": "system",
                        "content": "你是PPT结构分析专家。请分析这张PPT页面，判断它是什么类型。只返回JSON。"
                    },
                    {"role": "user", "content": prompt},
                ],
                response_format={"type": "json_object"},
                temperature=0.0,
            )
            result = response.choices[0].message.content.strip()
            import json
            data = json.loads(result)
            return data.get("page_type", "content")
        except Exception:
            # AI调用失败，默认归为内容页
            return "content"

    def _build_classification_prompt(self, index: int, text_areas: List[str], placeholders: List[Placeholder]) -> str:
        ph_info = [
            f"- 占位{idx+1}: 类型={p.type}, 位置(x={p.x}, y={p.y}), 大小(w={p.width}, h={p.height})"
            for idx, p in enumerate(placeholders)
        ]
        ph_text = "\n".join(ph_info)

        text_text = "\n".join([f"- {t}" for t in text_areas if t])

        return f"""请分析这张PPT幻灯片：

页面索引: {index}
现有文字内容:
{text_text}

占位区域:
{ph_text}

请判断这是什么类型的页面，可选类型:
- cover: 封面页
- toc: 目录页
- title: 标题页（只有一个大标题，没有内容）
- content: 内容页
- section: 章节分隔页
- blank: 空白页

请以JSON格式输出，格式:
{{
  "page_type": "你判断的类型"
}}
只输出JSON，不要其他内容。
"""
