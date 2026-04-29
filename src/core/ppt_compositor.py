# -*- coding: utf-8 -*-
"""
模块四：PPT Compositor - 内容模板合并生成
将整理好的内容填充到模板对应位置，保持原有样式，输出最终PPT
"""

import io
from io import BytesIO
from typing import List, Dict, Optional, Any

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE

from .template_parser import TemplateStructure, Placeholder, SlideInfo
from .content_organizer import ContentPlan, ContentPage, DataBlock, ImageRequest
from .image_handler import ImageResult


class PPTCompositor:
    """PPT合成器"""

    def __init__(self):
        pass

    def compose(
        self,
        template_file: BytesIO,
        template_structure: TemplateStructure,
        content_plan: ContentPlan,
        image_results: List[ImageResult],
    ) -> BytesIO:
        """合成最终PPT

        新逻辑：根据content_plan中的page type，从模板挑选对应类型的页面复制过来填充内容
        """

        prs_original = Presentation(template_file)

        # 获取模板库中每种类型的可用页面
        template_type_map: Dict[str, List[SlideInfo]] = {}
        for slide_info in template_structure.slides:
            if slide_info.page_type not in template_type_map:
                template_type_map[slide_info.page_type] = []
            template_type_map[slide_info.page_type].append(slide_info)

        # 保存原始presentation到output
        output = BytesIO()
        prs_original.save(output)
        output.seek(0)

        prs = Presentation(output)

        # 按content_plan填充内容
        for content_page in content_plan.pages:
            self._fill_content_page(prs, content_page, template_type_map)

        # 插入图片
        self._insert_images(prs, content_plan, image_results, template_type_map)

        # 输出
        output_final = BytesIO()
        prs.save(output_final)
        output_final.seek(0)
        return output_final

    def _fill_content_page(
        self,
        prs: Presentation,
        content_page: ContentPage,
        template_type_map: Dict[str, List[SlideInfo]],
    ):
        """填充一页内容，从模板找对应类型"""

        target_type = content_page.template_page_type

        # 获取模板中对应类型的页面，选第一个
        if target_type not in template_type_map or len(template_type_map[target_type]) == 0:
            # fallback to content type
            if "content" in template_type_map and len(template_type_map["content"]) > 0:
                target_type = "content"
            else:
                # if no matching, use any available
                if len(template_type_map) > 0:
                    target_type = next(iter(template_type_map.keys()))

        # pick the first template slide of this type
        template_slide_info = template_type_map[target_type][0]

        # get the actual slide from presentation by its original index
        target_slide = None
        for idx, slide in enumerate(prs.slides):
            if idx == template_slide_info.index:
                target_slide = slide
                break

        if target_slide is None:
            return

        self._fill_slide_content(target_slide, content_page, template_slide_info)

    def _fill_slide_content(self, slide, content_page: ContentPage, template_slide: SlideInfo):
        """填充页面文字内容，分配到多个占位符"""

        if not content_page.content or not content_page.content.strip():
            return

        # 找到这页的文字占位
        text_placeholders = [p for p in template_slide.placeholders if p.type == "text"]

        # 拆分内容，按空行分隔段落
        content_blocks = [block.strip() for block in content_page.content.split('\n\n') if block.strip()]

        if len(text_placeholders) == 0:
            return

        if len(text_placeholders) == 1:
            # 只有一个文字占位，放全部内容
            ph = text_placeholders[0]
            self._fill_text_placeholder(slide, ph, content_page.content)
        else:
            # 多个占位，拆分内容分配
            # 如果内容块比占位符少，剩下的保持原样
            for i, ph in enumerate(text_placeholders):
                if i < len(content_blocks):
                    self._fill_text_placeholder(slide, ph, content_blocks[i])
                else:
                    break

    def _fill_text_placeholder(self, slide, ph: Placeholder, content: str):
        """填充文字到占位符"""

        # 找到对应的形状
        for shape in slide.shapes:
            if shape.shape_id == ph.shape_id:
                if shape.has_text_frame:
                    # 清空原有内容
                    text_frame = shape.text_frame
                    text_frame.clear()
                    # 保留第一个段落，设置内容
                    p = text_frame.paragraphs[0]
                    p.text = content.strip()

                    # 自动调整字号
                    self._auto_fit_font(shape, content)
                break

    def _auto_fit_font(self, shape, content: str):
        """简单自动适配字号"""
        if not shape.has_text_frame:
            return

        text_frame = shape.text_frame
        # 估算合适的字号
        # 占位高度 (Emu -> Point)
        height_pt = shape.height / 12700
        lines = content.count("\n") + 1
        estimated_lines = max(lines, len(content) // 40 + 1)

        # 计算合适字号
        font_size = int(min(height_pt / estimated_lines * 0.8, 24))
        font_size = max(font_size, 10)  # 最小10号

        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(font_size)

    def _insert_images(
        self,
        prs: Presentation,
        content_plan: ContentPlan,
        image_results: List[ImageResult],
        template_type_map: Dict[str, List[SlideInfo]],
    ):
        """插入图片到对应位置"""

        for img_result in image_results:
            if not img_result.success or not img_result.image_data:
                continue

            # 找到对应类型的占位
            found_ph = None
            found_slide_info = None

            # 遍历模板找到对应类型和占位
            for slide_type, slides in template_type_map.items():
                if slide_type == img_result.template_page_type:
                    for slide_info in slides:
                        for ph in slide_info.placeholders:
                            if ph.type == img_result.placeholder_type:
                                found_ph = ph
                                found_slide_info = slide_info
                                break
                        if found_ph:
                            break
                    if found_ph:
                        break

            if not found_ph:
                continue

            # 找到对应原始slide in presentation
            found_slide = None
            for idx, slide in enumerate(prs.slides):
                if idx == found_slide_info.index:
                    found_slide = slide
                    break

            if not found_slide:
                continue

            # 转换位置（英寸 -> EMU）
            left = Emu(int(found_ph.x * 914400))
            top = Emu(int(found_ph.y * 914400))
            width = Emu(int(found_ph.width * 914400))
            height = Emu(int(found_ph.height * 914400))

            # 插入图片
            try:
                found_slide.shapes.add_picture(
                    BytesIO(img_result.image_data),
                    left,
                    top,
                    width,
                    height,
                )
            except Exception as e:
                print(f"Failed to insert image: {e}")

    def _create_charts(
        self,
        prs: Presentation,
        content_plan: ContentPlan,
        template_type_map: Dict[str, List[SlideInfo]],
    ):
        """创建图表"""

        for data_block in content_plan.data_blocks:
            if data_block.recommended_type == "table":
                self._create_table(prs, data_block, template_type_map)
            else:
                self._create_chart(prs, data_block, template_type_map)

    def _create_table(self, prs, data_block: DataBlock, template_type_map: Dict[str, List[SlideInfo]]):
        """创建表格"""
        # 简化实现：将数据作为文字填入占位
        # 完整实现需要解析结构化数据创建真实表格
        pass

    def _create_chart(self, prs, data_block: DataBlock, template_type_map: Dict[str, List[SlideInfo]]):
        """创建图表"""
        # 完整实现需要根据结构化数据创建对应类型的图表
        # MVP 阶段：数据已经由AI分析推荐好类型，开发阶段进一步完善
        pass

    def _get_chart_type(self, recommended_type: str) -> XL_CHART_TYPE:
        """获取PPT图表类型"""
        mapping = {
            "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "line": XL_CHART_TYPE.LINE,
            "pie": XL_CHART_TYPE.PIE,
            "scatter": XL_CHART_TYPE.XY_SCATTER,
        }
        return mapping.get(recommended_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
